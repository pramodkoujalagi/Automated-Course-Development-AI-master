from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import io

def ppt_generator(course_data, filename):
    
    # Create a new presentation
    presentation = Presentation()

    # Title Slide
    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title = title_slide.shapes.title
    title.text = course_data['title'].upper()  # Capitalize the title
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # Center align title
    title.text_frame.paragraphs[0].font.bold = True  # Bolden the title

    for module_name, module_data in course_data['content'].items():
        # Module Title and Content Slide
        module_slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        module_title = module_slide.shapes.title
        module_title.text = module_name
        module_title.text_frame.paragraphs[0].font.bold = True  # Set module title to bold
        module_title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # Center align module title
        module_body = module_slide.shapes.placeholders[1]
        module_body_text = "\n".join([lesson_name for lesson_name in module_data.keys()])
        module_body.text = module_body_text

        for lesson_name, lesson_data in module_data.items():
            # Lesson Title Slide
            lesson_title_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            placeholders = lesson_title_slide.shapes.placeholders
            if placeholders:
                slide_width = placeholders[0].width
                slide_height = placeholders[0].height
                left = (slide_width - Inches(8)) / 2  # Center the textbox
                top = (slide_height - Inches(1)) / 2  # Center the textbox
            else:
                slide_width, slide_height = Inches(10), Inches(7.5)  # Fallback dimensions
                left, top = Inches(1), Inches(3)  # Fallback position
            width, height = Inches(8), Inches(1)  # Adjust width and height
            lesson_title_box = lesson_title_slide.shapes.add_textbox(left, top, width, height)
            text_frame = lesson_title_box.text_frame
            text_frame.text = lesson_name
            text_frame.paragraphs[0].font.size = Pt(36)
            text_frame.paragraphs[0].font.bold = True  # Set lesson title to bold
            text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # Center align lesson title
            text_frame.word_wrap = True  # Enable text wrapping

            for slide_num, slide_content in lesson_data.items():
                # Split content into multiple slides if needed
               #  content_chunks = [slide_content[i:i+4] for i in range(0, len(slide_content), 4)] #To keep 4 pointers per slide
                content_chunks = [slide_content[i:i+3] for i in range(0, len(slide_content), 3)] #To keep 3 pointers per slide
                for idx, chunk in enumerate(content_chunks):
                    # Lesson Content Slide
                    content_slide = presentation.slides.add_slide(presentation.slide_layouts[1])
                    content_title = content_slide.shapes.title
                    content_title.text = slide_num + (" (Continued)" if idx > 0 else "")
                    content_title.text_frame.paragraphs[0].font.bold = True  # Set content title to bold
                    content_title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # Center align content title
                    content_body = content_slide.shapes.placeholders[1]
                    content_body.text = '\n'.join(chunk)
                    
                    # Set font size for content bullet points
                    content_body_frame = content_body.text_frame

                    for paragraph in content_body_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(25)

    # Thank You Slide
    thank_you_slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    # Add text box
    left = Inches(1)  # Adjust left position as needed
    slide_height = presentation.slide_height  # Get slide height directly
    text_height = Inches(1)  # Adjust text box height as needed
    top = (slide_height - text_height) / 2  # Position text box vertically centered
    width, height = Inches(8), Inches(1)  # Adjust width and height as needed
    thank_you_box = thank_you_slide.shapes.add_textbox(left, top, width, height)
 
    # Set text properties
    text_frame = thank_you_box.text_frame
    text_frame.text = "Keep Learning, Keep Growing!"
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.font.size = Pt(32)
    paragraph.font.bold = True
    text_frame.word_wrap = True

    course_name = course_data['title'].replace(" ", "")

    # presentation.save(f'{filename}.pptx')
    
    ppt_bytes = io.BytesIO()
    presentation.save(ppt_bytes)
    ppt_bytes.seek(0)

    return ppt_bytes

# if __name__ == "__main__":
#     course_data = {
#    "title":"Accounting and Finance",
#    "content":{
#       "Module 1: Introduction to Accounting and Finance":{
        #  "Lesson 1.1: Overview of Accounting and Finance":{
        #     "Overview of Accounting and Finance":[
        #        "Welcome to the World of Accounting and Finance!",
        #        "Accounting and finance are the backbone of any business, providing critical information for decision-making and strategy development.",
        #        "In this lesson, we'll embark on an exciting journey to explore the fundamentals of accounting and finance, laying the groundwork for a deeper understanding of the subject."
        #     ],
        #     "What is Accounting?":[
        #        "Definition: Accounting is the process of recording, classifying, and reporting financial transactions and events of a business.",
        #        "Objective: To provide stakeholders with accurate and timely financial information for informed decision-making."
        #     ],
        #     "What is Finance?":[
        #        "Definition: Finance is the management of money and investments for individuals, businesses, and organizations.",
        #        "Objective: To maximize wealth and minimize risk through optimal resource allocation."
        #     ],
        #     "Fundamental Accounting Principles":[
        #        "Accounting Entity: Separate accounting records for the business.",
        #        "Going Concern: Assumption that the business will continue to operate for the foreseeable future.",
        #        "Monetary Unit: Transactions are recorded in a common currency.",
        #        "Historical Cost: Assets and liabilities are recorded at their original cost.",
        #        "Matching Principle: Expenses are matched with revenues in the same period."
        #     ],
        #     "Overview of Financial Statements":[
        #        "Balance Sheet: Snapshot of assets, liabilities, and equity at a point in time.",
        #        "Income Statement: Summary of revenues and expenses over a period of time.",
        #        "Cash Flow Statement: Inflows and outflows of cash over a period of time."
        #     ],
        #     "Real-World Applications":[
        #        "Financial Analysis: Evaluate a company's performance using ratio analysis, trend analysis, and industry comparisons.",
        #        "Investment Decisions: Use financial statements to assess investment opportunities and risks."
        #     ],
        #     "Reflective Question":[
        #        "How do you think accounting and finance impact business decision-making?"
        #     ],
        #     "Exercise":[
        #        "Analyze a publicly traded company's financial statements to identify trends and areas for improvement."
        #     ],
        #     "Economics and Accounting":[
        #        "Microeconomics: Study of individual economic units, influencing accounting decisions.",
        #        "Macroeconomics: Study of the economy as a whole, impacting accounting policies."
        #     ],
        #     "Technology and Accounting":[
        #        "Accounting Information Systems (AIS): Use of technology to record, classify, and report financial transactions."
        #     ],
        #     "Recap: Overview of Accounting and Finance":[
        #        "Accounting: The process of recording, classifying, and reporting financial transactions and events.",
        #        "Finance: The management of money and investments for individuals, businesses, and organizations.",
        #        "Financial Statements: Balance Sheet, Income Statement, and Cash Flow Statement."
        #     ],
        #     "Key Takeaways":[
        #        "Understand the fundamental principles of accounting and finance.",
        #        "Appreciate the importance of financial statements in business decision-making.",
        #        "Recognize the interdisciplinary connections between accounting, finance, economics, and technology."
        #     ]
        #  },
        #  "Lesson 1.2: Financial Statements and Ratio Analysis":{
        #     "Unlocking the Secrets of Financial Statements and Ratio Analysis":[
        #        "Financial statements are the backbone of financial analysis, providing insights into a company's performance, profitability, and financial health."
        #     ],
        #     "What are Financial Statements?":[
        #        "Financial statements are formal records of a company's financial activities, providing stakeholders with a snapshot of its financial position and performance.",
        #        "Types of Financial Statements:",
        #        "Balance Sheet: A snapshot of the company's financial position at a specific point in time, showcasing its assets, liabilities, and equity.",
        #        "Income Statement: A summary of the company's revenues and expenses over a specific period, highlighting its profitability.",
        #        "Cash Flow Statement: A statement of the company's inflows and outflows of cash over a specific period."
        #     ],
        #     "Why are Financial Statements Important?":[
        #        "Financial statements help stakeholders:",
        #        "Evaluate a company's financial performance and position",
        #        "Identify areas of improvement and opportunity",
        #        "Make informed investment decisions",
        #        "Conduct peer comparisons and industry analysis"
        #     ],
        #     "Ratio Analysis: A Tool for Insight":[
        #        "Ratio analysis involves calculating and interpreting financial ratios to gain deeper insights into a company's performance.",
        #        "Types of Ratios:",
        #        "Liquidity Ratios: Measure a company's ability to pay its short-term debts (e.g., Current Ratio, Quick Ratio).",
        #        "Profitability Ratios: Measure a company's ability to generate earnings (e.g., Gross Margin Ratio, Return on Equity).",
        #        "Efficiency Ratios: Measure a company's ability to utilize its resources effectively (e.g., Asset Turnover Ratio, Inventory Turnover Ratio).",
        #        "Interpretation is Key"
        #     ],
        #     "The Balance Sheet: A Deeper Dive":[
        #        "Assets: Resources owned or controlled by the company (e.g., cash, inventory, property).",
        #        "Liabilities: Debts or obligations owed by the company (e.g., loans, accounts payable).",
        #        "Equity: The company's net worth or residual interest (e.g., common stock, retained earnings)."
        #     ],
        #     "The Income Statement: Unlocking Profitability":[
        #        "Revenues: Inflows from sales, services, or other activities (e.g., sales revenue, interest income).",
        #        "Expenses: Outflows or costs incurred to generate revenues (e.g., cost of goods sold, operating expenses).",
        #        "Net Income: The company's profitability after taxes and dividends."
        #     ],
        #     "Reflective Question 1.1":[
        #        "What are some potential limitations of relying solely on financial statements for investment decisions?"
        #     ],
        #     "Exercise 1.1: Financial Statement Analysis":[
        #        "Analyze a company's financial statements (e.g., Apple, Amazon) and calculate key ratios (e.g., current ratio, return on equity)."
        #     ],
        #     "Accounting Standards and Regulations":[
        #        "Financial statements are prepared in accordance with accounting standards like Generally Accepted Accounting Principles (GAAP) or International Financial Reporting Standards (IFRS)."
        #     ],
        #     "Interdisciplinary Connection:":[
        #        "Financial statements are used in various fields beyond accounting, such as finance, economics, and management."
        #     ],
        #     "Financial Statements and Ratio Analysis: A Recap":[
        #        "Financial statements provide insights into a company's financial position and performance.",
        #        "Ratio analysis is a valuable tool for evaluating a company's liquidity, profitability, and efficiency."
        #     ],
        #     "Key Takeaways:":[
        #        "Financial statements are essential for informed business decisions and investment analysis.",
        #        "Ratio analysis provides a deeper understanding of a company's financial performance.",
        #        "Critical thinking and contextual analysis are crucial for effective financial statement analysis."
        #     ]
        #  }
#       }
#    }
# }

#     ppt_generator(course_data)